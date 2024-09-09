import streamlit as st
import google.generativeai as genai
import os
from docx import Document
from docx.shared import Pt, RGBColor
from io import BytesIO
import matplotlib.pyplot as plt
from dotenv import load_dotenv
import json
import PyPDF2 as pdf
from pptx import Presentation

# Load environment variables
load_dotenv()

# Configure Gemini API
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Function to get response from Gemini API
def get_gemini_response(input_text):
    model = genai.GenerativeModel('gemini-pro')
    try:
        response = model.generate_content(input_text)
        if response and response.text:
            return response.text
        else:
            st.error("Received an empty response from the model.")
            return "{}"
    except Exception as e:
        st.error(f"Error while getting response from API: {str(e)}")
        return "{}"

# Function to extract text from uploaded PDF file
def input_pdf_text(uploaded_file):
    reader = pdf.PdfReader(uploaded_file)
    text = []
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text.append(page.extract_text())
    return text

# Function to extract text from uploaded Word document
def input_word_text(uploaded_file):
    doc = Document(uploaded_file)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return text

# Function to extract text from uploaded PPT file
def input_ppt_text(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return text

# Function to split text into manageable chunks
def split_text(text, max_chunk_size=2000):
    chunks = []
    while len(text) > max_chunk_size:
        chunk = text[:max_chunk_size]
        text = text[max_chunk_size:]
        chunks.append(chunk)
    chunks.append(text)
    return chunks

# Function to create an image for LaTeX formula
def create_formula_image(formula):
    plt.text(0.5, 0.5, f"${formula}$", fontsize=15, ha='center', va='center')
    plt.axis('off')

    image_stream = BytesIO()
    plt.savefig(image_stream, format='png', bbox_inches='tight', transparent=True)
    plt.close()
    image_stream.seek(0)
    return image_stream

# Function to save generated content with formula and formatted text to DOCX
def save_docx_with_formulas(generated_content):
    doc = Document()
    for content in generated_content:
        doc.add_heading(f"Page {content['Page']}", level=1)

        # Add Explanation with font settings
        explanation_paragraph = doc.add_paragraph(f"Explanation: {content['Explanation']}")
        explanation_paragraph.runs[0].font.size = Pt(12)
        explanation_paragraph.runs[0].font.name = 'Arial'
        explanation_paragraph.runs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Add Example
        doc.add_paragraph(f"Example: {content['Example']}")

        # Add Formula (as an image)
        formula_image = create_formula_image(content['Mini Test'])
        doc.add_paragraph("Mini Test (with formula):")
        doc.add_picture(formula_image, width=Pt(300))

        # Add Test Solution
        doc.add_paragraph(f"Test Solution: {content['Test Solution']}")

    file_stream = BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    return file_stream

# Expanded prompts
input_prompts = {
    "Mathematics": """
    You are an expert in mathematics. Your task is to explain the content on the given page in a detailed and comprehensive manner.
    Provide a relevant and contextual example to illustrate the concept, ensuring that your explanation and example are clear for someone without advanced knowledge.
    Lastly, create a mini-test related to the page content with at least two problems. Be sure to provide detailed solutions.
    Your response should be no less than 300-500 words.
    Page Content: {page_content}
    """,
    "Statistics": """
    You are an expert in statistics. Your task is to explain the content on the given page in a thorough and comprehensive manner.
    Break down statistical terms and concepts in simple language, and provide a real-world example.
    Additionally, create a mini-test with at least two problems and detailed solutions.
    Your response should be no less than 300-500 words.
    Page Content: {page_content}
    """,
    "Computer Science": """
    You are an expert in computer science. Your task is to explain the content on the given page in a detailed manner.
    Provide a practical example, and create a mini-test with at least two questions to test the reader's understanding.
    Include step-by-step solutions.
    Your response should be no less than 300-500 words.
    Page Content: {page_content}
    """
}

# Streamlit App
st.set_page_config(page_title="Interactify")
st.title("Interactify")

# Dropdown for subject selection
subject = st.selectbox("Select Subject", ["Mathematics", "Statistics", "Computer Science"])

# File uploader
uploaded_file = st.file_uploader("Upload Your Document (PDF, DOCX, PPTX, TXT)...", type=["pdf", "docx", "pptx", "txt"])

# Text input for page range
page_range_input = st.text_input("Enter page ranges (e.g., 78-79):")

# Submit button for processing the document
submit = st.button("Submit")

if submit:
    if uploaded_file:
        try:
            # Extract text
            if uploaded_file.type == "application/pdf":
                document_text = input_pdf_text(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                document_text = input_word_text(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                document_text = input_ppt_text(uploaded_file)
            elif uploaded_file.type == "text/plain":
                document_text = uploaded_file.read().decode("utf-8").split('\n')
            else:
                st.error("Unsupported file type!")
                st.stop()

            # Process page ranges
            page_ranges = []
            if page_range_input:
                try:
                    ranges = page_range_input.split(',')
                    for r in ranges:
                        start, end = map(int, r.split('-'))
                        page_ranges.append(range(start - 1, end))
                except ValueError:
                    st.error("Invalid page range format! Use 'start-end'.")
                    st.stop()
            else:
                page_ranges = [range(len(document_text))]

            # Generate content
            generated_content = []
            for range_set in page_ranges:
                for page_num in range_set:
                    if page_num < len(document_text):
                        page_content = document_text[page_num]
                        input_prompt = input_prompts[subject]
                        input_prompt_filled = input_prompt.format(page_content=page_content)
                        response = get_gemini_response(input_prompt_filled)
                        
                        try:
                            response_json = json.loads(response)
                            page_content_json = {
                                "Page": page_num + 1,
                                "Explanation": response_json.get("Explanation", "No explanation available."),
                                "Example": response_json.get("Example", "No example available."),
                                "Mini Test": response_json.get("Mini Test", "No mini test available."),
                                "Test Solution": response_json.get("Test Solution", "No test solution available.")
                            }
                            generated_content.append(page_content_json)
                        except json.JSONDecodeError:
                            st.error("Failed to decode JSON response.")
                    else:
                        st.warning(f"Page {page_num + 1} is out of range.")
            
            st.success("Content generated successfully.")
            
            # Allow download
            docx_file = save_docx_with_formulas(generated_content)
            st.download_button(
                label="Download Generated Content",
                data=docx_file,
                file_name="generated_content.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

        except Exception as e:
            st.error(f"Error: {str(e)}")
    else:
        st.warning("Please upload your document.")

# Collapsible section for questions
with st.expander("Ask a Question"):
    st.subheader("Ask Questions about Your Document")
    user_question = st.text_area("Type your question about the document content:")
    submit_question = st.button("Submit Question")
    
    if submit_question:
        if user_question:
            question_prompt = f"Based on the content of the document, answer the following question:\n{user_question}"
            response = get_gemini_response(question_prompt)
            st.markdown("### Answer to Your Question:")
            st.write(response)
        else:
            st.warning("Please type a question.")

